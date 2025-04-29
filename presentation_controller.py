"""
Presentation Controller System
=============================

A robust presentation control system that enables wireless control of PowerPoint/PDF presentations
using an ESP32 device via Bluetooth Low Energy (BLE).

This module implements a full-featured presentation management system with the following capabilities:
- Wireless presentation control using ESP32 via BLE
- Support for PowerPoint (pptx) and PDF presentations
- Real-time slide conversion and caching
- Fullscreen presentation mode
- Responsive GUI with presentation list and status indicators
- Automatic BLE reconnection handling
- Comprehensive logging system

Technical Specifications:
- Python 3.8+
- Dependencies: tkinter, bluepy, PyMuPDF, python-pptx, Pillow
- BLE Protocol: GATT (Generic Attribute Profile)
- Platform: Linux (with LibreOffice for conversion)

!! ALERT !!
In line 79, dont't forget to add your username or accordingly modify path suiting your project structure

"""

import tkinter as tk
from tkinter import ttk
import os
from bluepy import btle
import threading
import time
from pathlib import Path
import sys
import logging
from pptx import Presentation
import fitz  # PyMuPDF
import tempfile
import subprocess
from PIL import Image, ImageTk

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('/var/log/presentation/debug.log')
    ]
)
logger = logging.getLogger(__name__)

class PresentationSystem:
    def __init__(self):
        self.root = tk.Tk()
        self.root.title("Presentation Controller")
        
        # Get screen dimensions
        screen_width = self.root.winfo_screenwidth()
        screen_height = self.root.winfo_screenheight()
        logger.info(f"Screen dimensions: {screen_width}x{screen_height}")
        
        # Configure full screen
        self.root.attributes('-fullscreen', True)
        self.root.geometry(f"{screen_width}x{screen_height}+0+0")
        
        # Bind escape key to exit fullscreen
        self.root.bind('<Escape>', self.toggle_fullscreen)
        
        # BLE settings
        self.ESP32_MAC = # Your ESP32's MAC
        self.SERVICE_UUID = # Set your Service UUID
        self.CHARACTERISTIC_UUID = # Set your Characteristic UUID
        self.connected = False
        self.peripheral = None
        
        # Presentation variables
        self.presentations_dir = Path('/home/<user_name>/presentations')
        self.current_selection = 0
        self.presentation_mode = False
        self.current_slide = 0
        self.slides = []
        
        self.setup_gui()
        self.scan_presentations()
        
        # Start BLE connection in separate thread
        logger.info("Starting BLE thread")
        self.ble_thread = threading.Thread(target=self.maintain_ble_connection, daemon=True)
        self.ble_thread.start()
        
        # Add keyboard controls for testing
        self.setup_keyboard_controls()
        
        logger.info("Starting main loop")
        self.root.mainloop()

    def setup_gui(self):
        style = ttk.Style()
        style.configure('Large.TLabel', font=('Arial', 24))
        style.configure('Title.TLabel', font=('Arial', 32, 'bold'))
        
        self.root.configure(bg='black')
        self.main_frame = ttk.Frame(self.root, padding="20")
        self.main_frame.grid(row=0, column=0, sticky="nsew")
        
        self.root.grid_rowconfigure(0, weight=1)
        self.root.grid_columnconfigure(0, weight=1)
        
        self.title_label = ttk.Label(
            self.main_frame,
            text="Presentation Controller",
            style='Title.TLabel'
        )
        self.title_label.grid(row=0, column=0, pady=(0, 20))
        
        self.status_frame = ttk.Frame(self.main_frame)
        self.status_frame.grid(row=1, column=0, sticky="ew", pady=(0, 20))
        
        self.ble_status = ttk.Label(
            self.status_frame,
            text="⚫ Disconnected",
            style='Large.TLabel'
        )
        self.ble_status.grid(row=0, column=0, sticky="w")
        
        # Presentations List with larger font
        self.list_frame = ttk.LabelFrame(
            self.main_frame,
            text="Available Presentations",
            padding="10"
        )
        self.list_frame.grid(row=2, column=0, sticky="nsew")
        
        self.main_frame.grid_rowconfigure(2, weight=1)
        self.main_frame.grid_columnconfigure(0, weight=1)
        
        self.presentation_list = tk.Listbox(
            self.list_frame,
            font=('Arial', 20),
            bg='black',
            fg='white',
            selectmode=tk.SINGLE,
            selectbackground='#2c3e50',
            selectforeground='white'
        )
        self.presentation_list.pack(fill=tk.BOTH, expand=True)
        
        # Instructions
        instructions = """
        Controls:
        ↑ UP Button: Previous
        ↓ DOWN Button: Next
        Long Press: Select/Exit
        ESC: Toggle Fullscreen
        """
        self.instructions = ttk.Label(
            self.main_frame,
            text=instructions,
            style='Large.TLabel'
        )
        self.instructions.grid(row=3, column=0, pady=20)
        
        # Status bar
        self.status_bar = ttk.Label(
            self.main_frame,
            text="Ready",
            style='Large.TLabel'
        )
        self.status_bar.grid(row=4, column=0, sticky="ew")
        logger.info("GUI setup completed")

    def setup_keyboard_controls(self):
        self.root.bind('<Up>', lambda e: self.handle_command("UP"))
        self.root.bind('<Down>', lambda e: self.handle_command("DOWN"))
        self.root.bind('<Return>', lambda e: self.handle_command("SELECT"))
        logger.info("Keyboard controls set up")

    def toggle_fullscreen(self, event=None):
        is_fullscreen = not self.root.attributes('-fullscreen')
        self.root.attributes('-fullscreen', is_fullscreen)
        logger.info(f"Fullscreen toggled: {is_fullscreen}")

    def scan_presentations(self):
        logger.info("Scanning for presentations")
        self.presentation_list.delete(0, tk.END)
        try:
            presentations = sorted([
                f.name for f in self.presentations_dir.glob("*.ppt*")
            ])
            for pres in presentations:
                self.presentation_list.insert(tk.END, pres)
            
            if presentations:
                self.presentation_list.selection_set(0)
                self.status_bar.config(
                    text=f"Found {len(presentations)} presentations"
                )
                logger.info(f"Found {len(presentations)} presentations")
            else:
                self.status_bar.config(text="No presentations found")
                logger.info("No presentations found")
                
        except Exception as e:
            logger.error(f"Error scanning presentations: {e}")
            self.status_bar.config(text="Error scanning presentations")

    def maintain_ble_connection(self):
        max_retries = 3
        retry_delay = 2  # seconds
        
        while True:
            try:
                if not self.connected:
                    for retry in range(max_retries):
                        try:
                            logger.info(f"Connecting to {self.ESP32_MAC} (Attempt {retry + 1}/{max_retries})...")
                            # Try different address types
                            try:
                                self.peripheral = btle.Peripheral(self.ESP32_MAC, addrType=btle.ADDR_TYPE_RANDOM)
                            except:
                                self.peripheral = btle.Peripheral(self.ESP32_MAC, addrType=btle.ADDR_TYPE_PUBLIC)
                            
                            self.connected = True
                            logger.info("Connected!")
                            break
                        except Exception as e:
                            logger.error(f"Connection attempt {retry + 1} failed: {e}")
                            if retry < max_retries - 1:
                                time.sleep(retry_delay)
                            else:
                                raise

                    if self.connected:
                        self.root.after(0, lambda: self.ble_status.config(text="🟢 Connected"))
                        
                        time.sleep(0.5)
                        
                        # Notification handling
                        class NotifyDelegate(btle.DefaultDelegate):
                            def __init__(self, parent):
                                btle.DefaultDelegate.__init__(self)
                                self.parent = parent
                            
                            def handleNotification(self, cHandle, data):
                                command = data.decode()
                                logger.info(f"Received command: {command}")
                                # Schedule GUI updates in main thread
                                self.parent.root.after(0, lambda: 
                                    self.parent.handle_command(command))

                        self.peripheral.setDelegate(NotifyDelegate(self))
                        
                        try:
                            # Enable notifications with error handling
                            service = self.peripheral.getServiceByUUID(self.SERVICE_UUID)
                            characteristic = service.getCharacteristics(self.CHARACTERISTIC_UUID)[0]
                            
                            # Get the exact CCCD handle by searching all descriptors
                            cccd_handle = None
                            for desc in characteristic.getDescriptors():
                                logger.info(f"Found descriptor: {desc.uuid}")
                                if desc.uuid == btle.UUID("2902"): # for BLE CCCD
                                    cccd_handle = desc.handle
                                    logger.info(f"Found CCCD descriptor at handle: {cccd_handle}")
                                    break
                            
                            # Standard approach as fallback
                            if cccd_handle is None:
                                logger.info("CCCD not found, using calculated handle")
                                cccd_handle = characteristic.getHandle() + 1
                                
                            setup_data = b"\x01\x00"
                            
                            # Retry mechanism for enabling notifications
                            max_notify_retries = 3
                            for notify_retry in range(max_notify_retries):
                                try:
                                    logger.info(f"Enabling notifications on handle {cccd_handle}")
                                    self.peripheral.writeCharacteristic(cccd_handle, setup_data, withResponse=True)
                                    logger.info("Notifications enabled successfully")
                                    break
                                except Exception as notify_error:
                                    logger.error(f"Failed to enable notifications (attempt {notify_retry + 1}): {notify_error}")
                                    if notify_retry < max_notify_retries - 1:
                                        time.sleep(0.5)
                                    else:
                                        raise
                            
                            logger.info("Ready to receive notifications")
                            
                            # Keep connection alive and handle notifications
                            while self.connected:
                                try:
                                    if self.peripheral.waitForNotifications(1.0):
                                        continue
                                    time.sleep(0.1)
                                except Exception as notify_error:
                                    logger.error(f"Notification error: {notify_error}")
                                    raise
                        
                        except Exception as service_error:
                            logger.error(f"Service setup error: {service_error}")
                            raise
            except Exception as e:
                logger.error(f"BLE Error: {e}")
                self.connected = False
                self.root.after(0, lambda: 
                    self.ble_status.config(text="⚫ Disconnected"))
                
                # Clean up peripheral if it exists
                if self.peripheral:
                    try:
                        self.peripheral.disconnect()
                    except:
                        pass
                    self.peripheral = None
                
                time.sleep(2)

    def handle_command(self, command):
        """Handle commands from ESP32 or keyboard with debouncing"""
        logger.info(f"Received command: {command}")
        
        # Add debouncing for commands
        current_time = time.time()
        if hasattr(self, 'last_command_time') and current_time - self.last_command_time < 1.0:
            logger.info(f"Ignoring command due to debounce: {command}")
            return
        
        self.last_command_time = current_time
        
        if not self.presentation_mode:
            # Navigation mode
            if command == "UP":
                logger.info("Moving selection up")
                current = self.presentation_list.curselection()
                if current and current[0] > 0:
                    self.presentation_list.selection_clear(0, tk.END)
                    self.presentation_list.selection_set(current[0] - 1)
                    self.presentation_list.see(current[0] - 1)
                    logger.info(f"Selected item: {self.presentation_list.get(current[0] - 1)}")
                    
            elif command == "DOWN":
                logger.info("Moving selection down")
                current = self.presentation_list.curselection()
                if current and current[0] < self.presentation_list.size() - 1:
                    self.presentation_list.selection_clear(0, tk.END)
                    self.presentation_list.selection_set(current[0] + 1)
                    self.presentation_list.see(current[0] + 1)
                    logger.info(f"Selected item: {self.presentation_list.get(current[0] + 1)}")
                    
            elif command == "SELECT":
                logger.info("SELECT command received - starting presentation")
                self.start_presentation()
        else:
            # Presentation mode
            if command == "UP":
                logger.info("Showing previous slide")
                self.show_previous_slide()
            elif command == "DOWN":
                logger.info("Showing next slide")
                self.show_next_slide()
            elif command == "SELECT":
                logger.info("SELECT command received in presentation mode - ending presentation")
                self.end_presentation()

    def start_presentation(self):
        """Start the selected presentation"""
        current = self.presentation_list.curselection()
        if current:
            selected = self.presentation_list.get(current)
            logger.info(f"Starting presentation: {selected}")
            self.status_bar.config(text=f"Starting: {selected}")
            self.presentation_mode = True
            
            # Hide all UI elements
            self.main_frame.grid_remove()
            
            # Create presentation frame to take full space
            self.presentation_frame = ttk.Frame(self.root)
            self.presentation_frame.grid(row=0, column=0, sticky="nsew")
            
            # Create canvas for slides
            self.slide_canvas = tk.Canvas(
                self.presentation_frame,
                bg='black',
                highlightthickness=0
            )
            self.slide_canvas.pack(fill=tk.BOTH, expand=True)
            
            # Make sure root window grid is configured properly
            self.root.grid_rowconfigure(0, weight=1)
            self.root.grid_columnconfigure(0, weight=1)
            
            # Load presentation
            self.load_presentation(selected)

    def load_presentation(self, presentation_name):
        """Load the presentation and convert to images"""
        logger.info(f"Loading presentation: {presentation_name}")
        self.slides = []
        pptx_path = self.presentations_dir / presentation_name
        
        try:
            # Convert presentation to images
            with tempfile.TemporaryDirectory() as temp_dir:
                logger.info("Converting presentation to PDF")
                # Convert PPTX to PDF first
                output_pdf_name = presentation_name.replace('.pptx', '.pdf').replace('.ppt', '.pdf')
                pdf_path = os.path.join(temp_dir, output_pdf_name)
                
                # Run LibreOffice conversion with explicit output name
                subprocess.run([
                    'libreoffice',
                    '--headless',
                    '--convert-to',
                    'pdf',
                    '--outdir',
                    temp_dir,
                    str(pptx_path)
                ], check=True)
                
                if not os.path.exists(pdf_path):
                    raise FileNotFoundError(f"PDF conversion failed. Expected: {pdf_path}")
                
                logger.info("Converting PDF to images")
                # Convert PDF to images
                pdf_document = fitz.open(pdf_path)
                self.slide_images = []  # Store images in memory
                
                for page_num in range(pdf_document.page_count):
                    page = pdf_document[page_num]
                    pix = page.get_pixmap()
                    img_path = os.path.join(temp_dir, f"slide_{page_num}.png")
                    pix.save(img_path)
                    
                    # Load image into memory
                    img = Image.open(img_path)
                    self.slide_images.append(img.copy())
                    self.slides.append(page_num)  # Just store page number as reference
                
                logger.info(f"Converted {len(self.slides)} slides")
                
                # Force update the UI
                self.root.update_idletasks()
                
                # Show first slide immediately
                if self.slides:
                    self.current_slide = 0
                    self.show_current_slide()
                
        except Exception as e:
            logger.error(f"Error loading presentation: {e}")
            self.status_bar.config(text=f"Error loading presentation: {str(e)}")
            # Clean up on error
            self.end_presentation()

    def show_current_slide(self):
        """Display the current slide"""
        if 0 <= self.current_slide < len(self.slide_images):
            logger.info(f"Showing slide {self.current_slide + 1} of {len(self.slide_images)}")
            # Clear canvas
            self.slide_canvas.delete("all")
            
            try:
                # Get image from memory
                img = self.slide_images[self.current_slide]
                
                # Force update to ensure canvas has a valid size
                self.root.update_idletasks()
                
                # Scale image to fit canvas
                canvas_width = self.slide_canvas.winfo_width()
                canvas_height = self.slide_canvas.winfo_height()
                
                # Check for valid canvas size
                if canvas_width <= 0 or canvas_height <= 0:
                    logger.warning("Canvas not ready, using screen dimensions")
                    canvas_width = self.root.winfo_screenwidth()
                    canvas_height = self.root.winfo_screenheight()
                    
                if canvas_width > 0 and canvas_height > 0:
                    img_width, img_height = img.size
                    
                    scale = min(canvas_width/img_width, canvas_height/img_height)
                    new_width = int(img_width * scale)
                    new_height = int(img_height * scale)
                    
                    resized_img = img.resize((new_width, new_height), Image.LANCZOS)
                    self.photo = ImageTk.PhotoImage(resized_img)
                    
                    # Center image
                    x = (canvas_width - new_width) // 2
                    y = (canvas_height - new_height) // 2
                    
                    self.slide_canvas.create_image(
                        x, y, 
                        image=self.photo, 
                        anchor="nw"
                    )
                    
                    logger.info(f"Displayed slide with dimensions: {new_width}x{new_height}")
                    
                    # Update status if status bar exists
                    if hasattr(self, 'status_bar'):
                        self.status_bar.config(
                            text=f"Slide {self.current_slide + 1} of {len(self.slide_images)}"
                        )
                else:
                    logger.error("Canvas has invalid dimensions")
            except Exception as e:
                logger.error(f"Error displaying slide: {e}")
                if hasattr(self, 'status_bar'):
                    self.status_bar.config(text="Error displaying slide")

    def show_next_slide(self):
        """Show next slide"""
        if self.current_slide < len(self.slides) - 1:
            self.current_slide += 1
            self.show_current_slide()
            logger.info(f"Moved to next slide: {self.current_slide + 1}")

    def show_previous_slide(self):
        """Show previous slide"""
        if self.current_slide > 0:
            self.current_slide -= 1
            self.show_current_slide()
            logger.info(f"Moved to previous slide: {self.current_slide + 1}")

    def end_presentation(self):
        """End the current presentation"""
        logger.info("Ending presentation")
        self.presentation_mode = False
        
        # Remove presentation frame
        if hasattr(self, 'presentation_frame'):
            self.presentation_frame.destroy()
        
        # Restore main UI
        self.main_frame.grid()
        
        self.status_bar.config(text="Presentation ended")

if __name__ == "__main__":
    try:
        app = PresentationSystem()
        app.root.mainloop()
    except Exception as e:
        logger.error(f"Application error: {e}")
        sys.exit(1)